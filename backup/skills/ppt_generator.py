#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT 生成器 - 基于文字内容自动生成 PowerPoint 演示文稿
支持自定义模板、主题、布局
"""

import sys
import argparse
from pathlib import Path

if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8')

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
except ImportError:
    print("❌ 缺少依赖：python-pptx")
    print("请运行：pip install python-pptx")
    sys.exit(1)

# 主题配色方案
THEMES = {
    'default': {
        'primary': RGBColor(41, 98, 255),
        'secondary': RGBColor(52, 73, 94),
        'accent': RGBColor(231, 76, 60),
        'background': RGBColor(255, 255, 255),
        'text': RGBColor(44, 62, 80)
    },
    'dark': {
        'primary': RGBColor(52, 152, 219),
        'secondary': RGBColor(149, 165, 166),
        'accent': RGBColor(230, 126, 34),
        'background': RGBColor(44, 62, 80),
        'text': RGBColor(236, 240, 241)
    },
    'green': {
        'primary': RGBColor(39, 174, 96),
        'secondary': RGBColor(22, 160, 133),
        'accent': RGBColor(241, 196, 15),
        'background': RGBColor(255, 255, 255),
        'text': RGBColor(44, 62, 80)
    },
    'purple': {
        'primary': RGBColor(142, 68, 173),
        'secondary': RGBColor(155, 89, 182),
        'accent': RGBColor(243, 156, 18),
        'background': RGBColor(255, 255, 255),
        'text': RGBColor(44, 62, 80)
    }
}

def create_title_slide(prs, title, subtitle, theme):
    """创建标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme['primary']
    
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = subtitle
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = theme['secondary']
    
    return slide

def create_content_slide(prs, title, content_points, theme, slide_num):
    """创建内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme['primary']
    
    if content_points:
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        
        for i, point in enumerate(content_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(20)
            p.font.color.rgb = theme['text']
            p.space_after = Pt(14)
    
    page_num = slide.shapes.add_textbox(Inches(9), Inches(7), Inches(1), Inches(0.5))
    page_num.text_frame.text = f"{slide_num}"
    page_num.text_frame.paragraphs[0].font.size = Pt(12)
    page_num.text_frame.paragraphs[0].font.color.rgb = theme['secondary']
    
    return slide

def parse_content_text(text):
    """解析文本内容为幻灯片列表"""
    slides_content = []
    current_slide = {'title': '', 'points': []}
    
    for line in text.strip().split('\n'):
        line = line.strip()
        if not line:
            if current_slide['title'] or current_slide['points']:
                slides_content.append(current_slide)
                current_slide = {'title': '', 'points': []}
        elif not current_slide['title']:
            current_slide['title'] = line
        else:
            current_slide['points'].append(line)
    
    if current_slide['title'] or current_slide['points']:
        slides_content.append(current_slide)
    
    return slides_content

def generate_ppt(title, subtitle, content_text=None, output_file='output.pptx', theme_name='default'):
    """
    生成 PPT 演示文稿
    
    参数:
        title: 标题
        subtitle: 副标题
        content_text: 内容文本（可选）
        output_file: 输出文件路径
        theme_name: 主题名称
    
    返回:
        str: 输出文件路径
    """
    print(f'📊 开始生成 PPT...')
    print(f'标题：{title}')
    print(f'副标题：{subtitle}')
    print(f'主题：{theme_name}')
    print(f'输出：{output_file}')
    
    theme = THEMES.get(theme_name, THEMES['default'])
    prs = Presentation()
    
    # 创建标题页
    create_title_slide(prs, title, subtitle, theme)
    
    # 解析并创建内容页
    slides_content = []
    if content_text:
        slides_content = parse_content_text(content_text)
    
    if not slides_content:
        # 创建默认内容页
        slides_content = [{
            'title': '欢迎使用',
            'points': ['PPT 生成器', '支持多种主题', '简单易用']
        }]
    
    for i, slide_data in enumerate(slides_content, start=1):
        title = slide_data.get('title', f'第 {i} 页')
        points = slide_data.get('points', [])
        create_content_slide(prs, title, points, theme, i)
    
    prs.save(output_file)
    print(f'\n✅ PPT 生成成功！')
    print(f'共 {len(slides_content) + 1} 页')
    print(f'已保存：{output_file}')
    
    return output_file

def main():
    parser = argparse.ArgumentParser(description='PPT 生成器')
    parser.add_argument('--title', type=str, default='演示文稿', help='标题')
    parser.add_argument('--subtitle', type=str, default='AI 自动生成', help='副标题')
    parser.add_argument('--content', type=str, help='内容文本或文件路径')
    parser.add_argument('--output', type=str, default='output.pptx', help='输出文件路径')
    parser.add_argument('--theme', type=str, default='default',
                        choices=['default', 'dark', 'green', 'purple'],
                        help='主题')
    
    args = parser.parse_args()
    
    content_text = None
    if args.content:
        if Path(args.content).exists():
            content_text = Path(args.content).read_text(encoding='utf-8')
        else:
            content_text = args.content
    
    result = generate_ppt(
        title=args.title,
        subtitle=args.subtitle,
        content_text=content_text,
        output_file=args.output,
        theme_name=args.theme
    )
    
    sys.exit(0 if result else 1)

if __name__ == '__main__':
    main()
